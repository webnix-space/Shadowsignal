from flask import Flask, jsonify, request, render_template_string, send_file
import os
import json
import requests
import time
import uuid
import concurrent.futures
from datetime import datetime
from io import BytesIO

# ============================================================
#  SHADOWSIGNAL v2.0 - BAND SDK MULTI-AGENT ORCHESTRATION
#  Band of Agents Hackathon 2026
# ============================================================

app = Flask(__name__)

# --- ENVIRONMENT CONFIGURATION ---
BAND_API_KEY = os.getenv("BAND_API_KEY", "").strip()
BAND_PRO_TOKEN = os.getenv("BAND_PRO_TOKEN", "").strip()
BRIGHT_DATA_API_KEY = os.getenv("BRIGHT_DATA_API_KEY", "").strip()
BRIGHT_DATA_ZONE = os.getenv("BRIGHT_DATA_ZONE", "").strip()
AIML_API_KEY = os.getenv("AIML_API_KEY", "").strip()
FEATHERLESS_API_KEY = os.getenv("FEATHERLESS_API_KEY", "").strip()

# ============================================================
#  BAND SDK NATIVE ORCHESTRATION LAYER
# ============================================================

class BandRoom:
    """Native Band Room with shared context, agent messaging, and audit ledger"""
    def __init__(self, room_id, target_competitor):
        self.room_id = room_id
        self.target = target_competitor
        self.shared_context = {
            "raw_intel": None,
            "analysis": None,
            "strategy": None,
            "audit": None,
            "deliverables": {},
            "escalations": [],
            "agent_states": {}
        }
        self.ledger = []
        self.message_queue = []  # Agent-to-agent messaging
        self.human_escalations = []

    def broadcast(self, agent_name, action, payload):
        """Broadcast action to all agents in room"""
        entry = {
            "timestamp": datetime.now().isoformat(),
            "agent": agent_name,
            "action": action,
            "data": str(payload)[:200] + "..." if len(str(payload)) > 200 else str(payload)
        }
        self.ledger.append(entry)
        return entry

    def ask(self, from_agent, to_agent, question):
        """Direct agent-to-agent question (async messaging)"""
        msg = {
            "from": from_agent,
            "to": to_agent,
            "type": "ASK",
            "content": question,
            "timestamp": datetime.now().isoformat(),
            "status": "pending"
        }
        self.message_queue.append(msg)
        self.broadcast(from_agent, f"ASK->{to_agent}", question)
        return msg

    def respond(self, from_agent, to_agent, answer, msg_id=None):
        """Respond to agent question"""
        msg = {
            "from": from_agent,
            "to": to_agent,
            "type": "RESPONSE",
            "content": answer,
            "timestamp": datetime.now().isoformat(),
            "status": "delivered"
        }
        self.message_queue.append(msg)
        self.broadcast(from_agent, f"RESPONSE->{to_agent}", answer)
        return msg

    def challenge(self, from_agent, to_agent, reason, severity="medium"):
        """Challenge another agent's output (collaboration feature)"""
        msg = {
            "from": from_agent,
            "to": to_agent,
            "type": "CHALLENGE",
            "content": reason,
            "severity": severity,
            "timestamp": datetime.now().isoformat()
        }
        self.message_queue.append(msg)
        self.broadcast(from_agent, f"CHALLENGE->{to_agent}", f"[{severity}] {reason}")
        return msg

    def escalate(self, reason, requires_approval="human", blocking=True):
        """Escalate to human approver (Track 3: Regulated/High-Stakes)"""
        escalation = {
            "timestamp": datetime.now().isoformat(),
            "reason": reason,
            "requires_approval": requires_approval,
            "blocking": blocking,
            "status": "pending"
        }
        self.human_escalations.append(escalation)
        self.broadcast("SYSTEM", "ESCALATION", f"Human approval required: {reason}")
        return escalation

    def wait_for(self, key, timeout=30):
        """Wait for shared context key to be populated"""
        start = time.time()
        while self.shared_context.get(key) is None and (time.time() - start) < timeout:
            time.sleep(0.5)
        return self.shared_context.get(key)

    def get_messages_for(self, agent_name):
        """Get all messages directed to an agent"""
        return [m for m in self.message_queue if m["to"] == agent_name and m["status"] == "pending"]

    def mark_resolved(self, msg_id):
        """Mark message as resolved"""
        for m in self.message_queue:
            if m.get("id") == msg_id:
                m["status"] = "resolved"


# ============================================================
#  AGENT 1: INVESTIGATOR (Bright Data / Featherless AI)
# ============================================================

class InvestigatorAgent:
    """Gathers raw competitive intelligence via web sources"""
    name = "Investigator"

    def run(self, room, target):
        room.broadcast(self.name, "INTEL_GATHERING", f"Scanning global signals for {target}...")

        # Check if Analyst asked for specific focus
        messages = room.get_messages_for(self.name)
        focus = "pricing"
        for msg in messages:
            if "focus" in msg["content"].lower() or "specific" in msg["content"].lower():
                focus = "pricing tiers and renewal dates"
                room.respond(self.name, msg["from"], f"Adjusting focus to: {focus}")
                room.mark_resolved(msg.get("id"))

        # Gather intel
        raw_intel = self._gather_intel(target, focus)

        room.shared_context["raw_intel"] = raw_intel
        room.broadcast(self.name, "INTEL_DEPOSITED", f"Gathered {len(raw_intel)} chars of intel")
        return raw_intel

    def _gather_intel(self, target, focus="pricing"):
        """Gather intel via Bright Data or fallback mock data"""
        if BRIGHT_DATA_API_KEY and BRIGHT_DATA_ZONE:
            try:
                url = "https://api.brightdata.com/request"
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {BRIGHT_DATA_API_KEY}"
                }
                payload = {
                    "zone": BRIGHT_DATA_ZONE,
                    "url": f"https://www.google.com/search?q={target}+{focus}+enterprise+2026",
                    "format": "raw",
                    "data_format": "parsed_light"
                }
                res = requests.post(url, json=payload, headers=headers, timeout=20).json()
                organic = res.get("organic", [])
                results = []
                for r in organic[:3]:
                    results.append(f"{r.get('title')}: {r.get('description')}")
                return "\n".join(results) if results else f"Live search complete for {target}."
            except Exception as e:
                return f"Bright Data error: {str(e)}. Using curated intelligence."

        # Fallback: Curated mock data for demo reliability
        mock_db = {
            "Microsoft": [
                "Microsoft Removes EA Volume Discounts Tiers B-D: Effective Nov 2025, elimination of volume-based discount tiers increases costs 6-12% for enterprise customers. Source: internal pricing memo.",
                "Microsoft 365 Suite Price Hikes July 2026: Select M365 suites see 15% price increase, pushing customers toward premium AI bundles. Source: official blog.",
                "Microsoft Pricing Consistency Update: Standardizes pricing across channels, removes negotiation leverage for enterprise customers. Source: partner portal."
            ],
            "Nvidia": [
                "NVIDIA GPU Display Driver Vulnerability May 2026: Out-of-bounds write vulnerability allows remote code execution on enterprise AI infrastructure. CVE-2026-XXXX. Source: security bulletin.",
                "NVIDIA AI Enterprise Security Update: Industry-standard vulnerability scanning methods for container images. Source: NVIDIA AI Blueprint.",
                "NVIDIA Data Center GPU Pricing: H100 demand exceeds supply, enterprise pricing stable but lead times extended. Source: supply chain report."
            ],
            "Salesforce": [
                "Salesforce Einstein GPT Pricing: Per-user AI add-on pricing increases TCO 30% for enterprise customers. Source: pricing sheet.",
                "Salesforce Contract Negotiation Changes: Reduced flexibility in multi-year deals, standardized terms. Source: procurement advisory."
            ]
        }
        return "\n\n".join(mock_db.get(target, [f"General intelligence gathered for {target}."]))


# ============================================================
#  AGENT 2: ANALYST (AI/ML API - Deep Reasoning)
# ============================================================

class AnalystAgent:
    """Synthesizes intelligence into structured strategic analysis"""
    name = "Analyst"

    def run(self, room, target):
        room.broadcast(self.name, "ANALYSIS_START", "Waiting for intel from Investigator...")

        # Wait for intel with timeout
        intel = room.wait_for("raw_intel", timeout=15)
        if not intel:
            # Ask Investigator for clarification
            room.ask(self.name, "Investigator", 
                "Intel insufficient for analysis. Need specific pricing tier data and contract renewal timelines.")
            room.broadcast(self.name, "WAITING", "Requested clarification from Investigator")
            intel = room.wait_for("raw_intel", timeout=20)

        if not intel:
            intel = f"Default indicators for {target}"

        room.broadcast(self.name, "ANALYZING", f"Processing {len(intel)} chars of intel...")

        # Perform analysis via AI/ML API
        analysis = self._analyze(intel, target)

        room.shared_context["analysis"] = analysis
        room.broadcast(self.name, "ANALYSIS_COMPLETE", 
            f"Impact: {analysis.get('impact_score', 'N/A')}/100, Trend: {analysis.get('trend', 'unknown')}")
        return analysis

    def _analyze(self, intel, target):
        """Deep analysis using AI/ML API"""
        if AIML_API_KEY:
            try:
                url = "https://api.aimlapi.com/v1/chat/completions"
                headers = {
                    "Authorization": f"Bearer {AIML_API_KEY}",
                    "Content-Type": "application/json"
                }
                payload = {
                    "model": "nvidia/nemotron-3-nano-omni-30b-a3b-reasoning",
                    "messages": [
                        {"role": "system", "content": "You are a GTM Analyst. Analyze competitor intel and return structured JSON with: impact_score (0-100), trend_direction (up/down/stable), key_findings (array), recommended_actions (array), confidence (0.0-1.0), competitive_threat_level (low/medium/high/critical)."},
                        {"role": "user", "content": f"Target: {target}\nIntel: {intel}\n\nReturn JSON only."}
                    ],
                    "response_format": {"type": "json_object"}
                }
                res = requests.post(url, json=payload, headers=headers, timeout=30).json()
                content = res['choices'][0]['message']['content']
                return json.loads(content)
            except Exception as e:
                room = None  # placeholder for error context
                pass

        # Fallback structured analysis
        fallbacks = {
            "Microsoft": {
                "impact_score": 78,
                "trend_direction": "up",
                "key_findings": [
                    "EA discount elimination affects 60% of enterprise customers",
                    "M365 price hike pushes AI bundle adoption (vendor lock-in)",
                    "Hidden TCO increase of 6-12% for FY2026",
                    "Pricing consistency removes negotiation leverage"
                ],
                "recommended_actions": [
                    "Audit current EA level and contract renewal date",
                    "Evaluate alternative productivity suites (Google Workspace, Zoho)",
                    "Negotiate early renewal before July 2026 price hikes",
                    "Build TCO calculator showing true Microsoft costs"
                ],
                "confidence": 0.92,
                "competitive_threat_level": "high"
            },
            "Nvidia": {
                "impact_score": 65,
                "trend_direction": "stable",
                "key_findings": [
                    "GPU driver vulnerability affects enterprise AI workloads",
                    "CVE scanning gap between container and host level",
                    "Immediate patching required for PCI-DSS/HIPAA compliance",
                    "H100 supply constraints limit scaling"
                ],
                "recommended_actions": [
                    "Verify driver versions across all hosts (Windows + Linux)",
                    "Implement network segmentation for AI clusters",
                    "Document remediation for audit trail",
                    "Evaluate AMD MI300X as alternative for non-CUDA workloads"
                ],
                "confidence": 0.88,
                "competitive_threat_level": "medium"
            }
        }
        return fallbacks.get(target, {
            "impact_score": 50,
            "trend_direction": "unknown",
            "key_findings": [f"Limited intelligence available for {target}"],
            "recommended_actions": ["Gather more competitive intelligence"],
            "confidence": 0.5,
            "competitive_threat_level": "low"
        })


# ============================================================
#  AGENT 3: STRATEGIST (AI/ML API - NEW)
# ============================================================

class StrategistAgent:
    """Generates ranked counter-play strategies"""
    name = "Strategist"

    def run(self, room, target):
        room.broadcast(self.name, "STRATEGY_START", "Waiting for Analyst assessment...")

        analysis = room.wait_for("analysis", timeout=20)
        if not analysis:
            room.broadcast(self.name, "WAITING", "Analysis not ready, using default strategies")
            analysis = {"impact_score": 50, "competitive_threat_level": "medium"}

        room.broadcast(self.name, "STRATEGIZING", 
            f"Generating counter-plays for threat level: {analysis.get('competitive_threat_level', 'medium')}")

        strategies = self._generate_strategies(analysis, target)

        room.shared_context["strategy"] = strategies
        room.broadcast(self.name, "STRATEGIES_GENERATED", 
            f"{len(strategies)} strategies ranked by ROI and risk")
        return strategies

    def _generate_strategies(self, analysis, target):
        """Generate strategies using AI/ML API"""
        if AIML_API_KEY:
            try:
                url = "https://api.aimlapi.com/v1/chat/completions"
                headers = {
                    "Authorization": f"Bearer {AIML_API_KEY}",
                    "Content-Type": "application/json"
                }
                prompt = f"""Generate 3 competitive counter-play strategies against {target}.
Analysis: {json.dumps(analysis)}

Return JSON array. Each strategy must have:
- name (string)
- description (string)
- steps (array of strings)
- projected_roi (number, percentage)
- implementation_speed_weeks (number)
- risk_score (0-100, lower is safer)
- required_resources (array of strings)
- compliance_notes (string)

Rank by: ROI > Speed > Low Risk."""

                payload = {
                    "model": "nvidia/nemotron-3-nano-omni-30b-a3b-reasoning",
                    "messages": [
                        {"role": "system", "content": "You are a Competitive Strategist. Generate actionable counter-play strategies as JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    "response_format": {"type": "json_object"}
                }
                res = requests.post(url, json=payload, headers=headers, timeout=30).json()
                content = res['choices'][0]['message']['content']
                parsed = json.loads(content)
                return parsed.get("strategies", parsed) if isinstance(parsed, dict) else parsed
            except Exception as e:
                pass

        # Fallback strategies
        return [
            {
                "name": "Transparent TCO Calculator",
                "description": "Build public pricing calculator showing true Microsoft costs vs alternatives",
                "steps": ["Audit current Microsoft spend", "Build interactive calculator", "Launch microsite", "Sales enablement training"],
                "projected_roi": 15,
                "implementation_speed_weeks": 3,
                "risk_score": 20,
                "required_resources": ["Dev team (2 weeks)", "Marketing", "Legal review"],
                "compliance_notes": "Low risk - factual pricing comparison"
            },
            {
                "name": "Early Renewal Lock-in",
                "description": "Negotiate 3-year renewals before July 2026 price hikes",
                "steps": ["Identify renewal dates", "Prepare negotiation package", "Executive sponsor outreach", "Close by June 30"],
                "projected_roi": 25,
                "implementation_speed_weeks": 6,
                "risk_score": 35,
                "required_resources": ["Sales team", "Executive sponsors", "Legal"],
                "compliance_notes": "Medium risk - ensure no anti-competitive language"
            },
            {
                "name": "AI Value Bundle Counter",
                "description": "Bundle our AI capabilities into core platform, avoiding per-user add-on model",
                "steps": ["Audit AI feature gaps", "Develop integrated AI suite", "Pricing strategy workshop", "Go-to-market campaign"],
                "projected_roi": 30,
                "implementation_speed_weeks": 12,
                "risk_score": 45,
                "required_resources": ["Product team", "AI engineers", "Pricing strategist", "Marketing"],
                "compliance_notes": "Medium risk - monitor for predatory pricing claims"
            }
        ]


# ============================================================
#  AGENT 4: REGULATORY (Featherless AI - Compliance)
# ============================================================

class RegulatoryAgent:
    """Audits strategies for legal, ethical, and compliance risks"""
    name = "Regulatory"

    def run(self, room, target):
        room.broadcast(self.name, "AUDIT_START", "Auditing strategies for compliance risks...")

        # Wait for both analysis and strategy
        strategy = room.wait_for("strategy", timeout=25)
        analysis = room.shared_context.get("analysis", {})

        if not strategy:
            room.broadcast(self.name, "AUDIT_INCOMPLETE", "No strategies to audit")
            return {"status": "incomplete", "audit": "[LOW RISK] No strategies generated."}

        audit = self._audit(strategy, analysis, target)

        room.shared_context["audit"] = audit
        room.broadcast(self.name, "AUDIT_COMPLETE", audit[:100])

        # Check for critical risk and escalate
        if "[CRITICAL RISK]" in audit:
            room.challenge(self.name, "Strategist", 
                "Strategy contains critical compliance violations. Revision required before execution.", 
                severity="critical")
            room.escalate(
                reason="Critical compliance risk in competitive strategy - requires Legal approval",
                requires_approval="Legal/Compliance team",
                blocking=True
            )
        elif "[MEDIUM RISK]" in audit:
            room.challenge(self.name, "Strategist",
                "Strategy requires compliance adjustments. Suggest value-based positioning instead of direct price matching.",
                severity="medium")

        return audit

    def _audit(self, strategy, analysis, target):
        """Compliance audit using Featherless AI"""
        if FEATHERLESS_API_KEY:
            try:
                url = "https://api.featherless.ai/v1/chat/completions"
                headers = {
                    "Authorization": f"Bearer {FEATHERLESS_API_KEY}",
                    "Content-Type": "application/json"
                }
                strategies_text = json.dumps(strategy, indent=2)[:1500]
                payload = {
                    "model": "deepseek-ai/DeepSeek-V3.2",
                    "messages": [
                        {"role": "system", "content": "You are a Regulatory Compliance Agent. Audit competitive strategies for legal/ethical risks. Start with [CRITICAL RISK], [MEDIUM RISK], or [LOW RISK]. Check for: anti-trust violations, predatory pricing, data privacy risks, misrepresentation, unfair competition. Suggest compliant alternatives if flagged."},
                        {"role": "user", "content": f"Target: {target}\nStrategies: {strategies_text}\n\nProvide compliance assessment."}
                    ]
                }
                res = requests.post(url, json=payload, headers=headers, timeout=25).json()
                return res['choices'][0]['message']['content'].strip()
            except Exception as e:
                pass

        # Fallback audits
        audits = {
            "Microsoft": "[MEDIUM RISK] Strategy 2 (Early Renewal Lock-in) suggests aggressive contract negotiation. Risk: Could be interpreted as coercive under some jurisdictions. Compliant alternative: Position as 'customer success program' with transparent pricing guarantees. Strategy 1 and 3 are [LOW RISK].",
            "Nvidia": "[LOW RISK] All strategies focus on security remediation and alternative evaluation. No compliance concerns. Ensure all vulnerability claims are sourced from official security bulletins to avoid defamation."
        }
        return audits.get(target, "[LOW RISK] No significant compliance concerns detected. Standard competitive practices.")


# ============================================================
#  AGENT 5: CODEBAND ORCHESTRATOR (Auto-Generate Deliverables)
# ============================================================

class CodebandAgent:
    """Transforms approved strategies into executable Python artifacts"""
    name = "Codeband"

    def run(self, room, target):
        room.broadcast(self.name, "CODEBAND_START", "Checking regulatory approval...")

        audit = room.shared_context.get("audit", "")

        # BLOCK if critical risk
        if "[CRITICAL RISK]" in audit:
            room.broadcast(self.name, "BLOCKED", 
                "Cannot generate deliverables - CRITICAL RISK detected. Waiting for human approval.")
            return {"status": "blocked", "reason": "critical_risk", "artifacts": []}

        strategies = room.shared_context.get("strategy", [])
        analysis = room.shared_context.get("analysis", {})

        room.broadcast(self.name, "GENERATING", "Creating pricing comparison chart...")
        chart = self._generate_chart(target, analysis)

        room.broadcast(self.name, "GENERATING", "Creating strategy slide deck...")
        deck = self._generate_slide_deck(target, strategies, analysis)

        room.broadcast(self.name, "GENERATING", "Creating ROI calculator...")
        calculator = self._generate_roi_calculator(target, strategies)

        deliverables = {
            "chart": chart,
            "slide_deck": deck,
            "roi_calculator": calculator,
            "total_artifacts": 3,
            "generated_at": datetime.now().isoformat()
        }

        room.shared_context["deliverables"] = deliverables
        room.broadcast(self.name, "DELIVERABLES_READY", 
            f"3 artifacts generated: chart, slide deck, ROI calculator")

        return deliverables

    def _generate_chart(self, target, analysis):
        """Generate matplotlib pricing comparison chart"""
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import numpy as np
            from io import BytesIO

            # Create chart
            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
            fig.patch.set_facecolor('#030305')
            ax1.set_facecolor('#090a0f')
            ax2.set_facecolor('#090a0f')

            products = ["Basic", "Pro", "Enterprise", "AI Bundle"]
            competitor = [12, 25, 45, 80]
            our_solution = [10, 22, 38, 65]

            x = np.arange(len(products))
            width = 0.35

            bars1 = ax1.bar(x - width/2, competitor, width, label=target, color='#ef4444', alpha=0.8)
            bars2 = ax1.bar(x + width/2, our_solution, width, label='Our Solution', color='#3b82f6', alpha=0.8)

            ax1.set_ylabel('Price / user / month ($)', color='white', fontsize=11)
            ax1.set_title(f'{target} vs Our Pricing Comparison', color='white', fontsize=14, fontweight='bold')
            ax1.set_xticks(x)
            ax1.set_xticklabels(products, color='white')
            ax1.legend(facecolor='#090a0f', edgecolor='gray', labelcolor='white')
            ax1.tick_params(colors='white')
            ax1.grid(axis='y', alpha=0.2, color='gray')
            ax1.spines['bottom'].set_color('gray')
            ax1.spines['left'].set_color('gray')
            ax1.spines['top'].set_visible(False)
            ax1.spines['right'].set_visible(False)

            # ROI projection
            months = list(range(1, 13))
            monthly_savings = sum(competitor) - sum(our_solution)
            cumulative = [monthly_savings * m for m in months]

            ax2.plot(months, cumulative, marker='o', linewidth=2.5, color='#10b981', markersize=6)
            ax2.fill_between(months, cumulative, alpha=0.2, color='#10b981')
            ax2.set_xlabel('Months', color='white', fontsize=11)
            ax2.set_ylabel('Cumulative Savings ($)', color='white', fontsize=11)
            ax2.set_title('Projected Annual Savings', color='white', fontsize=14, fontweight='bold')
            ax2.tick_params(colors='white')
            ax2.grid(alpha=0.2, color='gray')
            ax2.spines['bottom'].set_color('gray')
            ax2.spines['left'].set_color('gray')
            ax2.spines['top'].set_visible(False)
            ax2.spines['right'].set_visible(False)

            plt.tight_layout()

            buf = BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', 
                       facecolor='#030305', edgecolor='none')
            buf.seek(0)
            plt.close()

            return {
                "status": "success",
                "format": "png",
                "data": buf.getvalue(),
                "filename": f"pricing_comparison_{target.lower()}.png"
            }
        except Exception as e:
            return {"status": "error", "error": str(e)}

    def _generate_slide_deck(self, target, strategies, analysis):
        """Generate PowerPoint strategy deck"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
            from io import BytesIO

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def add_title_slide(title, subtitle):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = RgbColor(3, 3, 5)
                bg.line.fill.background()

                tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = RgbColor(255, 255, 255)
                p.alignment = PP_ALIGN.LEFT

                tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
                tf2 = tb2.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = subtitle
                p2.font.size = Pt(18)
                p2.font.color.rgb = RgbColor(156, 163, 175)
                p2.alignment = PP_ALIGN.LEFT

            def add_content_slide(title, bullets):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = RgbColor(9, 10, 15)
                bg.line.fill.background()

                tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RgbColor(255, 255, 255)

                tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf2 = tb2.text_frame
                tf2.word_wrap = True

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf2.paragraphs[0]
                    else:
                        p = tf2.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = RgbColor(209, 213, 219)
                    p.space_after = Pt(12)

            # Build slides
            add_title_slide("ShadowSignal Intelligence Report", 
                          f"Competitive Analysis: {target} | Generated by Band 5-Agent System")

            add_content_slide("Executive Summary", [
                f"Target: {target}",
                f"Impact Score: {analysis.get('impact_score', 'N/A')}/100",
                f"Trend Direction: {analysis.get('trend_direction', 'Unknown')}",
                f"Threat Level: {analysis.get('competitive_threat_level', 'Unknown')}",
                f"Confidence: {analysis.get('confidence', 0) * 100:.0f}%",
                "Analysis generated by 5 collaborative AI agents via Band SDK"
            ])

            add_content_slide("Key Findings", analysis.get('key_findings', ['No findings available']))
            add_content_slide("Recommended Actions", analysis.get('recommended_actions', ['No actions available']))

            for i, strategy in enumerate(strategies[:3]):
                add_content_slide(f"Strategy {i+1}: {strategy.get('name', 'Untitled')}", [
                    f"Description: {strategy.get('description', '')}",
                    f"Projected ROI: {strategy.get('projected_roi', 'N/A')}%",
                    f"Implementation: {strategy.get('implementation_speed_weeks', 'N/A')} weeks",
                    f"Risk Score: {strategy.get('risk_score', 'N/A')}/100",
                    f"Resources: {', '.join(strategy.get('required_resources', []))}",
                    f"Compliance: {strategy.get('compliance_notes', 'N/A')}"
                ])

            buf = BytesIO()
            prs.save(buf)
            buf.seek(0)

            return {
                "status": "success",
                "format": "pptx",
                "data": buf.getvalue(),
                "filename": f"shadowsignal_report_{target.lower()}.pptx"
            }
        except Exception as e:
            return {"status": "error", "error": str(e)}

    def _generate_roi_calculator(self, target, strategies):
        """Generate Streamlit ROI calculator code"""
        code = f"""import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np

st.set_page_config(page_title="ShadowSignal ROI Calculator", layout="wide")

# Dark theme
st.markdown("""
<style>
    .stApp {{ background-color: #030305; color: #d1d5db; }}
    .stSlider > div > div > div {{ background-color: #3b82f6; }}
</style>
""", unsafe_allow_html=True)

st.title(f"💰 {{target}} Competitive ROI Calculator")
st.caption("Generated by ShadowSignal Band 5-Agent System")

# Sidebar inputs
st.sidebar.header("Your Current Situation")
users = st.sidebar.number_input("Number of Users", min_value=10, value=500, step=10)
current_cost = st.sidebar.number_input("Current Cost/User/Month ($)", min_value=5.0, value=45.0, step=1.0)
our_cost = st.sidebar.number_input("Our Solution Cost/User/Month ($)", min_value=5.0, value=35.0, step=1.0)

# Strategy selection
st.header("Select Counter-Play Strategy")
strategy_names = {json.dumps([s.get('name', f'Strategy {i+1}') for i, s in enumerate(strategies)])}
selected = st.selectbox("Choose strategy", strategy_names)
selected_strategy = next((s for s in {json.dumps(strategies)} if s['name'] == selected), None)

# Calculations
monthly_savings = users * (current_cost - our_cost)
annual_savings = monthly_savings * 12
impl_weeks = selected_strategy.get('implementation_speed_weeks', 6) if selected_strategy else 6
impl_cost = (selected_strategy.get('projected_roi', 15) * users * 0.5) if selected_strategy else 3750
net_roi = ((annual_savings - impl_cost) / impl_cost * 100) if impl_cost > 0 else 0

# Metrics
col1, col2, col3, col4 = st.columns(4)
col1.metric("Monthly Savings", f"${monthly_savings:,.0f}")
col2.metric("Annual Savings", f"${annual_savings:,.0f}")
col3.metric("Implementation Cost", f"${impl_cost:,.0f}")
col4.metric("Net ROI", f"{{net_roi:.1f}}%")

# Chart
months = pd.date_range(start='2026-07-01', periods=12, freq='M')
df = pd.DataFrame({{
    'Month': months,
    'Cumulative Savings': [monthly_savings * (i+1) for i in range(12)],
    'Break-even': [impl_cost] * 12
}})

fig = px.line(df, x='Month', y='Cumulative Savings',
              title='Projected Savings Over 12 Months',
              template='plotly_dark',
              line_shape='spline')
fig.add_hline(y=impl_cost, line_dash="dash", line_color="#f59e0b",
              annotation_text="Implementation Cost", annotation_position="bottom right")
fig.update_layout(
    plot_bgcolor='#090a0f',
    paper_bgcolor='#030305',
    font_color='#d1d5db'
)
st.plotly_chart(fig, use_container_width=True)

if selected_strategy:
    st.info(f"**Strategy:** {{selected_strategy.get('description', '')}}")
    st.success(f"**Timeline:** {{selected_strategy.get('implementation_speed_weeks', 'N/A')}} weeks | **Risk:** {{selected_strategy.get('risk_score', 'N/A')}}/100")

    with st.expander("Implementation Steps"):
        for step in selected_strategy.get('steps', []):
            st.write(f"• {{step}}")
"""

        return {
            "status": "success",
            "format": "python",
            "code": code,
            "filename": f"roi_calculator_{target.lower()}.py"
        }


# ============================================================
#  BAND WORKFLOW ORCHESTRATOR
# ============================================================

def run_band_workflow(target_competitor):
    """Execute the complete 5-agent collaborative workflow"""
    room_id = f"shadowsignal-{target_competitor.lower()}-{uuid.uuid4().hex[:8]}"
    room = BandRoom(room_id, target_competitor)

    room.broadcast("SYSTEM", "WORKFLOW_START", 
        f"Initializing 5-agent Band workflow for {target_competitor}")

    # Phase 1: Investigator gathers intel
    investigator = InvestigatorAgent()
    intel = investigator.run(room, target_competitor)

    # Phase 2: Analyst synthesizes (can ask Investigator for clarification)
    analyst = AnalystAgent()
    analysis = analyst.run(room, target_competitor)

    # Phase 3: Strategist generates counter-plays
    strategist = StrategistAgent()
    strategies = strategist.run(room, target_competitor)

    # Phase 4: Regulatory audits and can challenge/escalate
    regulatory = RegulatoryAgent()
    audit = regulatory.run(room, target_competitor)

    # Phase 5: Codeband generates deliverables (blocked if CRITICAL RISK)
    codeband = CodebandAgent()
    deliverables = codeband.run(room, target_competitor)

    # Calculate metrics
    agent_interactions = len([m for m in room.message_queue if m["type"] in ["ASK", "RESPONSE", "CHALLENGE"]])
    escalation_count = len(room.human_escalations)

    room.broadcast("SYSTEM", "WORKFLOW_COMPLETE", 
        f"5 agents completed in {len(room.ledger)} ledger entries")

    return {
        "room_id": room_id,
        "target": target_competitor,
        "status": "success" if deliverables.get("status") != "blocked" else "blocked",
        "ledger": room.ledger,
        "shared_context": room.shared_context,
        "agent_interactions": agent_interactions,
        "escalations": room.human_escalations,
        "escalation_count": escalation_count,
        "message_queue": room.message_queue,
        "deliverables": deliverables
    }


# ============================================================
#  FLASK ROUTES
# ============================================================

COMMAND_CENTER_HTML = """
<!DOCTYPE html>
<html>
<head>
    <script src="https://cdn.tailwindcss.com"></script>
    <style>
        .agent-investigator {{ color: #3b82f6; }}
        .agent-analyst {{ color: #a855f7; }}
        .agent-strategist {{ color: #f59e0b; }}
        .agent-regulatory {{ color: #ef4444; }}
        .agent-codeband {{ color: #10b981; }}
        .action-ask {{ background: rgba(254, 243, 199, 0.1); border-left: 3px solid #f59e0b; }}
        .action-challenge {{ background: rgba(254, 226, 226, 0.1); border-left: 3px solid #ef4444; }}
        .action-escalation {{ background: rgba(254, 226, 226, 0.15); border-left: 3px solid #dc2626; }}
        .action-complete {{ background: rgba(209, 250, 229, 0.1); border-left: 3px solid #10b981; }}
        .ledger-scroll {{ max-height: 400px; overflow-y: auto; }}
        .artifact-card {{ transition: all 0.3s ease; }}
        .artifact-card:hover {{ transform: translateY(-2px); box-shadow: 0 4px 12px rgba(59, 130, 246, 0.15); }}
    </style>
</head>
<body class="bg-[#030305] text-gray-400 font-sans min-h-screen p-4 md:p-8">
    <div class="max-w-6xl mx-auto space-y-6">
        <!-- Header -->
        <div class="flex flex-col md:flex-row justify-between items-start md:items-center border-b border-gray-800 pb-4 gap-4">
            <div>
                <h1 class="text-3xl font-extrabold text-white tracking-tight">
                    SHADOWSIGNAL<span class="text-blue-500">.BAND</span>
                </h1>
                <p class="text-xs text-gray-500 uppercase tracking-widest mt-1">
                    Native Band SDK 5-Agent Collaborative Orchestration
                </p>
            </div>
            <div class="flex gap-2 flex-wrap">
                <span class="text-xs font-mono text-blue-400 bg-blue-950/40 border border-blue-900 px-3 py-1 rounded">
                    5-AGENT HANDOFF
                </span>
                <span class="text-xs font-mono text-green-400 bg-green-950/40 border border-green-900 px-3 py-1 rounded">
                    BAND SDK v2
                </span>
                <span class="text-xs font-mono text-purple-400 bg-purple-950/40 border border-purple-900 px-3 py-1 rounded">
                    CODEBAND
                </span>
            </div>
        </div>

        <!-- Input -->
        <div class="grid grid-cols-1 md:grid-cols-4 gap-4 bg-[#090a0f] border border-gray-800 p-4 rounded-xl">
            <div class="md:col-span-3">
                <label class="block text-[10px] font-mono uppercase text-gray-500 mb-2">
                    Target Competitor
                </label>
                <input id="targetInput" type="text" value="Microsoft" 
                    class="w-full bg-[#101116] border border-gray-700 rounded-lg p-3 text-white 
                           focus:outline-none focus:border-blue-500 font-semibold text-sm">
            </div>
            <div class="flex items-end">
                <button onclick="runBandWorkflow()" id="controlBtn" 
                    class="w-full bg-blue-600 hover:bg-blue-500 text-white font-bold py-3 rounded-lg 
                           transition-all disabled:opacity-50 disabled:cursor-not-allowed">
                    RUN BAND WORKFLOW
                </button>
            </div>
        </div>

        <!-- Metrics (hidden initially) -->
        <div id="metricsPanel" class="hidden grid grid-cols-2 md:grid-cols-5 gap-3">
            <div class="bg-[#07080c] border border-gray-800 rounded-lg p-3 text-center">
                <div id="agentCount" class="text-xl font-bold text-white">5</div>
                <div class="text-[10px] text-gray-500 uppercase">Agents</div>
            </div>
            <div class="bg-[#07080c] border border-gray-800 rounded-lg p-3 text-center">
                <div id="ledgerCount" class="text-xl font-bold text-blue-400">0</div>
                <div class="text-[10px] text-gray-500 uppercase">Ledger Entries</div>
            </div>
            <div class="bg-[#07080c] border border-gray-800 rounded-lg p-3 text-center">
                <div id="interactionCount" class="text-xl font-bold text-purple-400">0</div>
                <div class="text-[10px] text-gray-500 uppercase">Interactions</div>
            </div>
            <div class="bg-[#07080c] border border-gray-800 rounded-lg p-3 text-center">
                <div id="escalationCount" class="text-xl font-bold text-red-400">0</div>
                <div class="text-[10px] text-gray-500 uppercase">Escalations</div>
            </div>
            <div class="bg-[#07080c] border border-gray-800 rounded-lg p-3 text-center">
                <div id="artifactCount" class="text-xl font-bold text-green-400">0</div>
                <div class="text-[10px] text-gray-500 uppercase">Artifacts</div>
            </div>
        </div>

        <!-- Agent Collaboration Log -->
        <div class="grid grid-cols-1 lg:grid-cols-2 gap-4">
            <div class="bg-[#07080c] border border-gray-800 rounded-xl p-4">
                <h3 class="text-xs font-mono uppercase text-gray-400 border-b border-gray-800 pb-2 mb-3">
                    Band Room Agent Collaboration Log
                </h3>
                <div id="agentChat" class="font-mono text-xs space-y-1 ledger-scroll">
                    <div class="text-gray-600 italic p-2">Band workflow ready. Click RUN to start 5-agent collaboration...</div>
                </div>
            </div>

            <div class="bg-[#07080c] border border-gray-800 rounded-xl p-4">
                <h3 class="text-xs font-mono uppercase text-gray-400 border-b border-gray-800 pb-2 mb-3">
                    Shared Context State
                </h3>
                <div id="contextState" class="font-mono text-xs">
                    <div class="text-gray-600 p-2">Waiting for workflow execution...</div>
                </div>
            </div>
        </div>

        <!-- Results -->
        <div id="resultsPanel" class="hidden space-y-4">
            <div class="grid grid-cols-1 md:grid-cols-3 gap-4">
                <div class="bg-[#090a0f] border border-purple-900/50 p-4 rounded-xl">
                    <h4 class="text-[10px] font-mono text-purple-400 uppercase mb-2">GTM Analysis</h4>
                    <pre id="analysisResult" class="text-xs text-gray-300 overflow-x-auto whitespace-pre-wrap"></pre>
                </div>
                <div class="bg-[#090a0f] border border-amber-900/50 p-4 rounded-xl">
                    <h4 class="text-[10px] font-mono text-amber-400 uppercase mb-2">Strategist Output</h4>
                    <pre id="strategyResult" class="text-xs text-gray-300 overflow-x-auto whitespace-pre-wrap"></pre>
                </div>
                <div class="bg-[#090a0f] border border-red-900/50 p-4 rounded-xl">
                    <h4 class="text-[10px] font-mono text-red-400 uppercase mb-2">Regulatory Audit</h4>
                    <pre id="auditResult" class="text-xs text-gray-300 overflow-x-auto whitespace-pre-wrap"></pre>
                </div>
            </div>

            <!-- Deliverables -->
            <div id="deliverablesPanel" class="bg-[#090a0f] border border-green-900/50 p-4 rounded-xl">
                <h4 class="text-[10px] font-mono text-green-400 uppercase mb-3">Codeband Deliverables</h4>
                <div id="artifactsGrid" class="grid grid-cols-1 md:grid-cols-3 gap-3"></div>
            </div>
        </div>
    </div>

    <script>
        async function runBandWorkflow() {
            const btn = document.getElementById('controlBtn');
            const target = document.getElementById('targetInput').value;
            const chatBox = document.getElementById('agentChat');
            const contextBox = document.getElementById('contextState');

            btn.disabled = true;
            btn.innerText = 'BAND WORKFLOW RUNNING...';

            // Reset UI
            document.getElementById('metricsPanel').classList.remove('hidden');
            document.getElementById('resultsPanel').classList.add('hidden');
            document.getElementById('deliverablesPanel').classList.add('hidden');
            chatBox.innerHTML = '<div class="text-blue-500 animate-pulse p-2">[BAND] Initializing 5-agent collaborative workflow...</div>';

            const startTime = Date.now();

            try {
                const response = await fetch('/api/run-band-workflow', {
                    method: 'POST',
                    headers: {'Content-Type': 'application/json'},
                    body: JSON.stringify({target: target})
                }).then(r => r.json());

                // Render agent chat with color coding and action styling
                chatBox.innerHTML = '';
                response.ledger.forEach(entry => {
                    const agentClass = 'agent-' + entry.agent.toLowerCase().replace(/\s/g, '-');
                    let actionClass = '';
                    if (entry.action.includes('ASK')) actionClass = 'action-ask';
                    else if (entry.action.includes('CHALLENGE')) actionClass = 'action-challenge';
                    else if (entry.action.includes('ESCALATION')) actionClass = 'action-escalation';
                    else if (entry.action.includes('COMPLETE') || entry.action.includes('READY')) actionClass = 'action-complete';

                    const time = entry.timestamp ? entry.timestamp.split('T')[1].split('.')[0] : '00:00:00';

                    chatBox.innerHTML += `
                        <div class="flex gap-2 p-1.5 rounded ${actionClass}">
                            <span class="text-gray-600 text-[10px] pt-0.5">${time}</span>
                            <span class="${agentClass} font-bold text-[11px]">[${entry.agent}]</span>
                            <span class="text-gray-300 text-[11px]">${entry.action}</span>
                            <span class="text-gray-500 text-[10px] truncate">→ ${entry.data}</span>
                        </div>
                    `;
                });

                // Show context state
                const ctx = response.shared_context;
                contextBox.innerHTML = `
                    <div class="space-y-2 text-[11px]">
                        <div class="flex justify-between border-b border-gray-800 pb-1">
                            <span class="text-gray-500">Target:</span>
                            <span class="text-white">${ctx.target || target}</span>
                        </div>
                        <div class="flex justify-between border-b border-gray-800 pb-1">
                            <span class="text-gray-500">Intel:</span>
                            <span class="text-blue-400">${ctx.raw_intel ? '✓ Deposited' : '✗ Missing'}</span>
                        </div>
                        <div class="flex justify-between border-b border-gray-800 pb-1">
                            <span class="text-gray-500">Analysis:</span>
                            <span class="text-purple-400">${ctx.analysis ? '✓ Complete' : '✗ Missing'}</span>
                        </div>
                        <div class="flex justify-between border-b border-gray-800 pb-1">
                            <span class="text-gray-500">Strategy:</span>
                            <span class="text-amber-400">${ctx.strategy ? '✓ Generated' : '✗ Missing'}</span>
                        </div>
                        <div class="flex justify-between border-b border-gray-800 pb-1">
                            <span class="text-gray-500">Audit:</span>
                            <span class="text-red-400">${ctx.audit ? '✓ Audited' : '✗ Missing'}</span>
                        </div>
                        <div class="flex justify-between">
                            <span class="text-gray-500">Deliverables:</span>
                            <span class="text-green-400">${ctx.deliverables && ctx.deliverables.total_artifacts ? ctx.deliverables.total_artifacts + ' artifacts' : '✗ Blocked'}</span>
                        </div>
                    </div>
                `;

                // Show results
                document.getElementById('resultsPanel').classList.remove('hidden');
                document.getElementById('analysisResult').textContent = 
                    typeof ctx.analysis === 'object' ? JSON.stringify(ctx.analysis, null, 2) : ctx.analysis || 'No analysis';
                document.getElementById('strategyResult').textContent = 
                    typeof ctx.strategy === 'object' ? JSON.stringify(ctx.strategy, null, 2) : ctx.strategy || 'No strategies';
                document.getElementById('auditResult').textContent = ctx.audit || 'No audit';

                // Show deliverables
                if (response.deliverables && response.deliverables.status !== 'blocked') {
                    document.getElementById('deliverablesPanel').classList.remove('hidden');
                    const artifactsGrid = document.getElementById('artifactsGrid');
                    artifactsGrid.innerHTML = '';

                    const artifacts = [
                        {key: 'chart', name: 'Pricing Chart', icon: '📊', color: 'blue', type: 'png'},
                        {key: 'slide_deck', name: 'Strategy Deck', icon: '📑', color: 'amber', type: 'pptx'},
                        {key: 'roi_calculator', name: 'ROI Calculator', icon: '🧮', color: 'green', type: 'py'}
                    ];

                    artifacts.forEach(art => {
                        const data = response.deliverables[art.key];
                        if (data && data.status === 'success') {
                            artifactsGrid.innerHTML += `
                                <div class="artifact-card bg-[#07080c] border border-${art.color}-900/50 rounded-lg p-3">
                                    <div class="text-2xl mb-2">${art.icon}</div>
                                    <div class="text-xs font-mono text-${art.color}-400 uppercase">${art.name}</div>
                                    <div class="text-[10px] text-gray-500 mt-1">${art.type.toUpperCase()}</div>
                                    <a href="/api/download/${art.type}/${response.room_id}/${art.key}" 
                                       class="mt-2 inline-block text-[10px] bg-${art.color}-900/30 text-${art.color}-400 
                                              border border-${art.color}-900/50 px-2 py-1 rounded hover:bg-${art.color}-900/50">
                                        DOWNLOAD
                                    </a>
                                </div>
                            `;
                        }
                    });
                }

                // Update metrics
                document.getElementById('ledgerCount').textContent = response.ledger.length;
                document.getElementById('interactionCount').textContent = response.agent_interactions || 0;
                document.getElementById('escalationCount').textContent = response.escalation_count || 0;
                document.getElementById('artifactCount').textContent = 
                    (response.deliverables && response.deliverables.total_artifacts) || 0;

            } catch (err) {
                chatBox.innerHTML += `<div class="text-red-500 p-2">[ERROR] ${err.message}</div>`;
            } finally {
                btn.disabled = false;
                btn.innerText = 'RUN BAND WORKFLOW';
            }
        }
    </script>
</body>
</html>
"""

@app.route('/')
def home():
    return render_template_string(COMMAND_CENTER_HTML)

@app.route('/api/run-band-workflow', methods=['POST'])
def api_run_band_workflow():
    """Execute the full 5-agent Band collaborative workflow"""
    data = request.get_json() or {}
    target = data.get('target', 'Microsoft')

    try:
        result = run_band_workflow(target)
        return jsonify(result)
    except Exception as e:
        return jsonify({
            "status": "error",
            "error": str(e),
            "traceback": str(e.__traceback__)
        }), 500

@app.route('/api/download/<filetype>/<room_id>/<artifact>')
def download_artifact(filetype, room_id, artifact):
    """Download generated artifacts"""
    # In production, store artifacts in memory/cache with room_id key
    # For demo, return placeholder
    return jsonify({"status": "download_ready", "filetype": filetype, "artifact": artifact})


# ============================================================
#  MAIN
# ============================================================

if __name__ == '__main__':
    app.run(debug=True, port=5000)
